# -*- coding: utf-8 -*-
"""
Functions for creating PowerPoint presentations from SDynPy objects.

This modules includes functions for assembling a PowerPoint presentation from
SDynPy objects, saving users the tediousness of putting a large number of images
and tables into the presentation.

Copyright 2022 National Technology & Engineering Solutions of Sandia,
LLC (NTESS). Under the terms of Contract DE-NA0003525 with NTESS, the U.S.
Government retains certain rights in this software.

This program is free software: you can redistribute it and/or modify
it under the terms of the GNU General Public License as published by
the Free Software Foundation, either version 3 of the License, or
(at your option) any later version.

This program is distributed in the hope that it will be useful,
but WITHOUT ANY WARRANTY; without even the implied warranty of
MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
GNU General Public License for more details.

You should have received a copy of the GNU General Public License
along with this program.  If not, see <https://www.gnu.org/licenses/>.
"""

import pptx
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE
import tempfile
import io
from PIL import Image
import matplotlib.pyplot as plt
import numpy as np
import os
import time

from ..core.sdynpy_shape import mac, matrix_plot


def position_placeholder(presentation, placeholder, left=None, top=None, right=None, bottom=None):
    sw = presentation.slide_width
    sh = presentation.slide_height
    if left is None:
        left = placeholder.left
    if left < 0:
        left = sw - abs(left)
    if right is None:
        right = placeholder.left + placeholder.width
    if right < 0:
        right = sw - abs(right)
    if top is None:
        top = placeholder.top
    if top < 0:
        top = sh - abs(top)
    if bottom is None:
        bottom = placeholder.top + placeholder.height
    if bottom < 0:
        bottom = sh - abs(bottom)
    width = right - left
    height = bottom - top
    placeholder.left = left
    placeholder.top = top
    placeholder.width = width
    placeholder.height = height


def add_title_slide(presentation, title, subtitle='', title_slide_layout_index=0):
    title_slide_layout = presentation.slide_layouts[title_slide_layout_index]
    slide = presentation.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    position_placeholder(presentation, slide.shapes.title, right=-slide.shapes.title.left)
    slide.placeholders[1].text = subtitle


def add_section_header_slide(presentation, title, subtitle='', section_header_slide_layout_index=2):
    section_header_slide_layout = presentation.slide_layouts[section_header_slide_layout_index]
    slide = presentation.slides.add_slide(section_header_slide_layout)
    slide.shapes.title.text = title
    position_placeholder(presentation, slide.shapes.title, right=-slide.shapes.title.left)
    slide.placeholders[1].text = subtitle


def add_geometry_overview_slide(presentation, geometry, title='Geometry',
                                geometry_plot_kwargs={},
                                animate_geometry=True,
                                save_animation_kwargs={'frames': 200, 'frame_rate': 20},
                                content_slide_layout_index=1):
    bullet_slide = presentation.slide_layouts[content_slide_layout_index]
    geometry_plot = geometry.plot(**geometry_plot_kwargs)[0]
    slide = presentation.slides.add_slide(bullet_slide)
    text_placeholder = slide.placeholders[1]
    if animate_geometry:
        output_save_file = os.path.join(tempfile.gettempdir(),
                                        'Geometry.gif')
        geometry_plot.save_rotation_animation(output_save_file, **save_animation_kwargs)
        geometry_plot.close()
        pic = slide.shapes.add_picture(output_save_file, Inches(1), Inches(1))
    else:
        with io.BytesIO() as output:
            time.sleep(0.5)
            img = Image.fromarray(geometry_plot.screenshot())
            img.save(output, format='PNG')
            pic = slide.shapes.add_picture(output, Inches(1), Inches(1))
    ar = pic.width / pic.height
    slide.shapes.title.text = title
    # Move the bullet points to the left
    position_placeholder(presentation, slide.placeholders[1], right=(
        presentation.slide_width - slide.placeholders[1].left) // 2)
    # Pic to the right
    position_placeholder(presentation, pic, left=(presentation.slide_width + slide.placeholders[1].left) // 2,
                         right=-slide.placeholders[1].left,
                         top=slide.placeholders[1].top)
    pic.height = int(pic.width / ar)
    geometry_plot.close()

    text_placeholder.text_frame.paragraphs[0].text = 'Geometry Information:'
    for label, array in zip(['Nodes', 'Coordinate Systems', 'Tracelines', 'Elements'],
                            [geometry.node, geometry.coordinate_system, geometry.traceline, geometry.element]):
        p = text_placeholder.text_frame.add_paragraph()
        p.text = '{:}: {:}'.format(label, array.size)
        p.level = 1


def add_shape_overview_slide(presentation, shapes, title='Modal Parameters',
                             exp_data=None, fit_data=None,
                             subplots_kwargs={}, plot_kwargs={}, axes_modifiers={},
                             mac_subplots_kwargs={}, matrix_plot_kwargs={},
                             frequency_format='{:0.2f}', damping_format='{:.03f}',
                             empty_slide_layout_index=5):
    empty_slide = presentation.slide_layouts[empty_slide_layout_index]
    slide = presentation.slides.add_slide(empty_slide)
    slide.shapes.title.text = title
    # Add a table for the mode shapes
    table_shape = slide.shapes.add_table(shapes.size + 1, 4, slide.shapes.title.left,
                                         slide.shapes.title.top +
                                         slide.shapes.title.height + Inches(0.5),
                                         presentation.slide_width // 2 - int(1.5 * slide.shapes.title.left), Inches(1))
    table = table_shape.table
    table.cell(0, 0).text = 'Mode'
    table.cell(0, 1).text = 'Freq (Hz)'
    table.cell(0, 2).text = 'Damp (%)'
    table.cell(0, 3).text = 'Description'

    table.columns[0].width = Inches(1)
    table.columns[1].width = Inches(1)
    table.columns[2].width = Inches(1)
    table.columns[3].width = presentation.slide_width // 2 - \
        int(1.5 * slide.shapes.title.left) - sum([table.columns[i].width for i in range(3)])

    for i, shape in enumerate(shapes.flatten()):
        table.cell(i + 1, 0).text = '{:}'.format(i + 1)
        table.cell(i + 1, 1).text = frequency_format.format(shape.frequency)
        table.cell(i + 1, 2).text = damping_format.format(shape.damping * 100)
        table.cell(i + 1, 3).text = shape.comment1

    mac_matrix = mac(shapes)
    fig, ax = plt.subplots(1, 1, **mac_subplots_kwargs)
    matrix_plot(mac_matrix, ax, **matrix_plot_kwargs)

    with io.BytesIO() as output:
        fig.savefig(output)
        pic = slide.shapes.add_picture(output, Inches(1), Inches(1))
        ar = pic.width / pic.height
        position_placeholder(presentation, pic, left=(presentation.slide_width + table_shape.left) // 2,
                             right=-table_shape.left,
                             top=table_shape.top,
                             bottom=table_shape.top + (presentation.slide_height - table_shape.top) // 2)
        pic.width = int(pic.height * ar)
        plt.close(fig)

    if not exp_data is None:
        fig, ax = plt.subplots(1, 1, **subplots_kwargs)
        h1 = ax.plot(exp_data.abscissa.T, exp_data.ordinate.T, 'r', **plot_kwargs)
        if not fit_data is None:
            h2 = ax.plot(fit_data.abscissa.T, fit_data.ordinate.T, 'b', **plot_kwargs)
            ax.legend([np.atleast_1d(h1)[0], np.atleast_1d(h2)[0]],
                      ['Experiment', 'Analytic Fit'])
        for key, val in axes_modifiers.items():
            getattr(ax, key)(val)
        if ('set_yscale', 'log'.lower()) in [(key, val) for key, val in axes_modifiers.items()]:
            ax.set_ylim(exp_data.ordinate.min() / 1.5, exp_data.ordinate.max() * 1.5)
        else:
            rng = exp_data.ordinate.max() - exp_data.ordinate.min()
            ax.set_ylim(exp_data.ordinate.min() - rng / 20, exp_data.ordinate.max() + rng / 20)
        with io.BytesIO() as output:
            fig.savefig(output)
            pic = slide.shapes.add_picture(output, Inches(1), Inches(1))
            ar = pic.width / pic.height
            plt.close(fig)
        position_placeholder(presentation, pic, left=(presentation.slide_width + table_shape.left) // 2,
                             right=-table_shape.left,
                             top=table_shape.top +
                             (presentation.slide_height - table_shape.top) // 2,
                             bottom=-slide.shapes.title.top)
        pic.width = int(pic.height * ar)


def add_shape_animation_slides(presentation, geometry, shapes, title_format='Mode {number:}',
                               text_format='Mode {number:}\nFrequency: {frequency:0.2f}\nDamping: {damping:0.3f}%',
                               save_animation_kwargs={'frames': 20, 'frame_rate': 20},
                               geometry_plot_kwargs={},
                               plot_shape_kwargs={},
                               content_slide_layout_index=1,
                               left_right=True):
    bullet_slide = presentation.slide_layouts[content_slide_layout_index]
    output_save_file = os.path.join(tempfile.gettempdir(),
                                    'Shape_{:}.gif')
    shapes_plot = geometry.plot_shape(
        shapes.flatten(), plot_kwargs=geometry_plot_kwargs, **plot_shape_kwargs)
    shapes_plot.save_animation_all_shapes(output_save_file, **save_animation_kwargs)
    shapes_plot.close()
    for i, shape in enumerate(shapes.flatten()):
        slide = presentation.slides.add_slide(bullet_slide)
        text_placeholder = slide.placeholders[1]
        pic = slide.shapes.add_picture(output_save_file.format(i + 1), Inches(1), Inches(1))
        ar = pic.width / pic.height
        slide.shapes.title.text = title_format.format(index=i, number=i + 1, frequency=shape.frequency,
                                                      damping=shape.damping * 100, modal_mass=shape.modal_mass,
                                                      comment1=shape.comment1, comment2=shape.comment2,
                                                      comment3=shape.comment3, comment4=shape.comment4,
                                                      comment5=shape.comment5)
        if left_right:
            # Move the bullet points to the left
            position_placeholder(presentation, text_placeholder,
                                 right=presentation.slide_width // 2 - Inches(0.25))
            # Pic to the right
            position_placeholder(presentation, pic, left=presentation.slide_width // 2 + Inches(0.25),
                                 right=-text_placeholder.left,
                                 top=text_placeholder.top)
        else:
            # Move the bullet points to the left
            position_placeholder(presentation, text_placeholder, right=-text_placeholder.left,
                                 bottom=text_placeholder.top + Inches(1))
            # Pic to the right
            position_placeholder(presentation, pic, left=text_placeholder.left,
                                 right=-text_placeholder.left,
                                 top=text_placeholder.top + text_placeholder.height + Inches(0.25))
        pic.height = int(pic.width / ar)
        # Check to see if it is off the screen
        distance_over = (pic.height + pic.top - presentation.slide_height +
                         Inches(0.25)) / pic.height
        if distance_over > 0:
            pic.height = int(pic.height * (1 - distance_over))
            pic.width = int(pic.width * (1 - distance_over))

        os.remove(output_save_file.format(i + 1))
        # Now add text
        for j, text_line in enumerate(text_format.split('\n')):
            text = text_line.replace('\t', '').format(index=i, number=i + 1, frequency=shape.frequency,
                                                      damping=shape.damping * 100, modal_mass=shape.modal_mass,
                                                      comment1=shape.comment1, comment2=shape.comment2,
                                                      comment3=shape.comment3, comment4=shape.comment4,
                                                      comment5=shape.comment5)
            if j == 0:
                paragraph = text_placeholder.text_frame.paragraphs[0]
            else:
                paragraph = text_placeholder.text_frame.add_paragraph()
            paragraph.text = text
            paragraph.level = text_line.count('\t')
        text_placeholder.text_frame.fit_text()
        text_placeholder.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def add_shape_comparison_overview_slide(presentation, shapes_1, shapes_2,
                                        title='Modal Parameter Comparison',
                                        shapes_1_label=None, shapes_2_label=None,
                                        subplots_kwargs={}, plot_kwargs={}, axes_modifiers={},
                                        frequency_format='{:0.2f}', damping_format='{:.03f}',
                                        mac_format='{:0.2f}', mac_subplots_kwargs={},
                                        matrix_plot_kwargs={},
                                        table_threshold=0.75,
                                        empty_slide_layout_index=5):
    empty_slide = presentation.slide_layouts[empty_slide_layout_index]
    slide = presentation.slides.add_slide(empty_slide)
    slide.shapes.title.text = title
    mac_matrix = mac(shapes_1, shapes_2)
    shape_indices = np.where(mac_matrix > table_threshold)
    # Plot the MAC matrix
    fig, ax = plt.subplots(1, 1, **mac_subplots_kwargs)
    matrix_plot(mac_matrix, ax, **matrix_plot_kwargs)
    if not shapes_1_label is None:
        ax.set_ylabel(shapes_1_label)
    if not shapes_2_label is None:
        ax.set_xlabel(shapes_2_label)
    fig.tight_layout()
    # Add a table for the mode shapes
    table_shape = slide.shapes.add_table(shape_indices[0].size + 1, 7, slide.shapes.title.left,
                                         slide.shapes.title.top +
                                         slide.shapes.title.height + Inches(0.5),
                                         presentation.slide_width // 2 - int(1.5 * slide.shapes.title.left), Inches(1))
    table = table_shape.table
    table.cell(0, 0).text = '{:} Mode'.format(
        'Shape 1' if shapes_1_label is None else shapes_1_label)
    table.cell(0, 1).text = 'Freq (Hz)'
    table.cell(0, 2).text = 'Damp (%)'
    table.cell(0, 3).text = '{:} Mode'.format(
        'Shape 2' if shapes_2_label is None else shapes_2_label)
    table.cell(0, 4).text = 'Freq (Hz)'
    table.cell(0, 5).text = 'Damp (%)'
    table.cell(0, 6).text = 'MAC'

    for i, indices in enumerate(zip(*shape_indices)):
        index_1, index_2 = indices
        shape_1 = shapes_1[index_1]
        shape_2 = shapes_2[index_2]
        table.cell(i + 1, 0).text = '{:}'.format(index_1 + 1)
        table.cell(i + 1, 3).text = '{:}'.format(index_2 + 1)
        table.cell(i + 1, 1).text = frequency_format.format(shape_1.frequency)
        table.cell(i + 1, 4).text = frequency_format.format(shape_2.frequency)
        table.cell(i + 1, 2).text = damping_format.format(shape_1.damping * 100)
        table.cell(i + 1, 5).text = damping_format.format(shape_2.damping * 100)
        table.cell(i + 1, 6).text = mac_format.format(mac_matrix[index_1, index_2])

    with io.BytesIO() as output:
        fig.savefig(output)
        pic = slide.shapes.add_picture(output, Inches(1), Inches(1))
        ar = pic.width / pic.height
        position_placeholder(presentation, pic, left=(presentation.slide_width + table_shape.left) // 2,
                             right=-table_shape.left,
                             top=table_shape.top,
                             bottom=-slide.shapes.title.top)
        pic.height = int(pic.width / ar)
        plt.close(fig)


def add_shape_comparison_animation_slides(presentation, geometry_1, shapes_1,
                                          geometry_2, shapes_2,
                                          title_format='Comparison for Mode {number:}',
                                          text_format='Mode {number:}\nFrequency: {frequency:0.2f}\nDamping: {damping:0.3f}%',
                                          save_animation_kwargs={'frames': 20, 'frame_rate': 20},
                                          geometry_plot_kwargs_1={},
                                          plot_shape_kwargs_1={},
                                          geometry_plot_kwargs_2={},
                                          plot_shape_kwargs_2={},
                                          two_content_slide_layout_index=3):
    two_content_slide = presentation.slide_layouts[two_content_slide_layout_index]
    output_save_file = os.path.join(tempfile.gettempdir(),
                                    'Shape_{:}.gif')
    slides = []
    for j, (shapes, geometry, geometry_plot_kwargs, plot_shape_kwargs) in enumerate(
            [(shapes_1, geometry_1, geometry_plot_kwargs_1, plot_shape_kwargs_1),
             (shapes_2, geometry_2, geometry_plot_kwargs_2, plot_shape_kwargs_2)]):
        shapes_plot = geometry.plot_shape(
            shapes.flatten(), plot_kwargs=geometry_plot_kwargs, **plot_shape_kwargs)
        shapes_plot.save_animation_all_shapes(output_save_file, **save_animation_kwargs)
        shapes_plot.close()
        for i, shape in enumerate(shapes.flatten()):
            if j == 0:
                slide = presentation.slides.add_slide(two_content_slide)
                slides.append(slide)
            else:
                slide = slides[i]
            text_placeholder = slide.shapes[1 + j]
            pic = slide.shapes.add_picture(output_save_file.format(i + 1), Inches(1), Inches(1))
            ar = pic.width / pic.height
            if j == 0:
                slide.shapes.title.text = title_format.format(index=i, number=i + 1, frequency=shape.frequency,
                                                              damping=shape.damping * 100, modal_mass=shape.modal_mass,
                                                              comment1=shape.comment1, comment2=shape.comment2,
                                                              comment3=shape.comment3, comment4=shape.comment4,
                                                              comment5=shape.comment5)
            position_placeholder(presentation, pic, left=text_placeholder.left,
                                 right=text_placeholder.left + text_placeholder.width,
                                 top=text_placeholder.top + Inches(1.5))
            pic.height = int(pic.width / ar)
            # Check to see if it is off the screen
            distance_over = (pic.height + pic.top - presentation.slide_height +
                             Inches(0.25)) / pic.height
            if distance_over > 0:
                pic.height = int(pic.height * (1 - distance_over))
                pic.width = int(pic.width * (1 - distance_over))

            os.remove(output_save_file.format(i + 1))
            # Now add text
            for k, text_line in enumerate(text_format.split('\n')):
                text = text_line.replace('\t', '').format(index=i, number=i + 1, frequency=shape.frequency,
                                                          damping=shape.damping * 100, modal_mass=shape.modal_mass,
                                                          comment1=shape.comment1, comment2=shape.comment2,
                                                          comment3=shape.comment3, comment4=shape.comment4,
                                                          comment5=shape.comment5)
                if k == 0:
                    paragraph = text_placeholder.text_frame.paragraphs[0]
                else:
                    paragraph = text_placeholder.text_frame.add_paragraph()
                paragraph.text = text
                paragraph.level = text_line.count('\t')
            text_placeholder.text_frame.fit_text()
            text_placeholder.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def create_summary_pptx(presentation, title=None, subtitle='',
                        geometry=None, shapes=None, frfs=None,
                        slide_width=None, slide_height=None,
                        max_shapes=None, max_frequency=None,
                        frequency_format='{:0.1f}', damping_format='{:.02f}',
                        cmif_kwargs={'part': 'imag', 'tracking': None},
                        subplots_kwargs={}, plot_kwargs={},
                        save_animation_kwargs={'frames': 20, 'frame_rate': 20},
                        geometry_plot_kwargs={},
                        title_slide_layout_index=0,
                        content_slide_layout_index=1,
                        empty_slide_layout_index=5
                        ):
    experimental_cmif = None if frfs is None else frfs.compute_cmif(**cmif_kwargs)
    frequencies = None if experimental_cmif is None else experimental_cmif[0].abscissa

    analytic_frfs = None if (shapes is None or frfs is None) else shapes.compute_frf(frequencies, np.unique(frfs.coordinate[..., 0]),
                                                                                     np.unique(frfs.coordinate[..., 1]))
    analytic_cmif = analytic_frfs.compute_cmif(**cmif_kwargs)

    if isinstance(presentation, pptx.presentation.Presentation):
        prs = presentation
    else:
        prs = pptx.Presentation()
        if not slide_width is None:
            prs.slide_width = Inches(slide_width)
        if not slide_height is None:
            prs.slide_height = Inches(slide_height)

    def position_placeholder(placeholder, left=None, top=None, right=None, bottom=None):
        sw = prs.slide_width
        sh = prs.slide_height
        if left is None:
            left = placeholder.left
        if left < 0:
            left = sw - abs(left)
        if right is None:
            right = placeholder.left + placeholder.width
        if right < 0:
            right = sw - abs(right)
        if top is None:
            top = placeholder.top
        if top < 0:
            top = sh - abs(top)
        if bottom is None:
            bottom = placeholder.top + placeholder.height
        if bottom < 0:
            bottom = sh - abs(bottom)
        width = right - left
        height = bottom - top
        placeholder.left = left
        placeholder.top = top
        placeholder.width = width
        placeholder.height = height

    # Add Title Slide
    if not title is None:
        title_slide_layout = prs.slide_layouts[title_slide_layout_index]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        position_placeholder(slide.shapes.title, right=-slide.shapes.title.left)
        slide.placeholders[1].text = subtitle
        position_placeholder(slide.placeholders[1], right=-slide.placeholders[1].left)

    bullet_slide = prs.slide_layouts[content_slide_layout_index]
    empty_slide = prs.slide_layouts[empty_slide_layout_index]

    # Plot the test geometry
    if not geometry is None:
        geometry_plot = geometry.plot(**geometry_plot_kwargs)[0]
        slide = prs.slides.add_slide(bullet_slide)
        text_placeholder = slide.placeholders[1]
        with io.BytesIO() as output:
            time.sleep(0.5)
            img = Image.fromarray(geometry_plot.screenshot())
            img.save(output, format='PNG')
            pic = slide.shapes.add_picture(output, Inches(1), Inches(1))
        slide.shapes.title.text = 'Test Geometry'
        # Move the bullet points to the left
        position_placeholder(slide.placeholders[1], right=(
            prs.slide_width - slide.placeholders[1].left) // 2)
        # Pic to the right
        position_placeholder(pic, left=(prs.slide_width + slide.placeholders[1].left) // 2,
                             right=-slide.placeholders[1].left,
                             top=slide.placeholders[1].top)
        pic.height = int(pic.width * img.size[1] / img.size[0])
        geometry_plot.close()

        text_placeholder.text_frame.paragraphs[0].text = 'Geometry Information:'
        for label, array in zip(['Nodes', 'Coordinate Systems', 'Tracelines', 'Elements'],
                                [geometry.node, geometry.coordinate_system, geometry.traceline, geometry.element]):
            p = text_placeholder.text_frame.add_paragraph()
            p.text = '{:}: {:}'.format(label, array.size)
            p.level = 1

    # Plot the list of shapes
    if not shapes is None:
        if not max_shapes is None:
            shapes = shapes.flatten()[:max_shapes]
        if not max_frequency is None:
            shapes = shapes.flatten()[shapes.flatten().frequency < max_frequency]
        slide = prs.slides.add_slide(empty_slide)

        slide.shapes.title.text = 'Modal Parameters'
        # Add a table for the mode shapes
        table_shape = slide.shapes.add_table(shapes.size + 1, 4, slide.shapes.title.left,
                                             slide.shapes.title.top +
                                             slide.shapes.title.height + Inches(0.5),
                                             prs.slide_width // 2 - int(1.5 * slide.shapes.title.left), Inches(1))
        table = table_shape.table
        table.cell(0, 0).text = 'Mode'
        table.cell(0, 1).text = 'Freq (Hz)'
        table.cell(0, 2).text = 'Damp (%)'
        table.cell(0, 3).text = 'Description'

        table.columns[0].width = Inches(1)
        table.columns[1].width = Inches(1)
        table.columns[2].width = Inches(1)
        table.columns[3].width = prs.slide_width // 2 - \
            int(1.5 * slide.shapes.title.left) - sum([table.columns[i].width for i in range(3)])

        for i, shape in enumerate(shapes.flatten()):
            table.cell(i + 1, 0).text = '{:}'.format(i + 1)
            table.cell(i + 1, 1).text = frequency_format.format(shape.frequency)
            table.cell(i + 1, 2).text = damping_format.format(shape.damping * 100)
            table.cell(i + 1, 3).text = shape.comment1

        if not analytic_cmif is None:
            fig, ax = plt.subplots(1, 1, **subplots_kwargs)
            experimental_cmif.plot(ax, **plot_kwargs)
            analytic_cmif.plot(ax, **plot_kwargs)
            ax.legend(['Experiment', 'Analytic Fit'])
            ax.set_yscale('log')
            ax.set_ylim(experimental_cmif.ordinate.min() / 1.5,
                        experimental_cmif.ordinate.max() * 1.5)
            with io.BytesIO() as output:
                fig.savefig(output)
                pic = slide.shapes.add_picture(output, Inches(1), Inches(1))
                ar = pic.width / pic.height
                plt.close(fig)
            position_placeholder(pic, left=(prs.slide_width + table_shape.left) // 2,
                                 right=-table_shape.left,
                                 top=table_shape.top)
            pic.height = int(pic.width / ar)

        # Now we need to plot each shape if possible
        if not geometry is None:
            output_save_file = os.path.join(tempfile.gettempdir(),
                                            'Shape_{:}.gif')
            shapes_plot = geometry.plot_shape(shapes.flatten(), plot_kwargs=geometry_plot_kwargs)
            shapes_plot.save_animation_all_shapes(output_save_file, **save_animation_kwargs)
            shapes_plot.close()
            for i, shape in enumerate(shapes.flatten()):
                slide = prs.slides.add_slide(bullet_slide)
                text_placeholder = slide.placeholders[1]
                pic = slide.shapes.add_picture(output_save_file.format(i + 1), Inches(1), Inches(1))
                ar = pic.width / pic.height
                slide.shapes.title.text = 'Mode {:}'.format(i + 1)
                # Move the bullet points to the left
                position_placeholder(text_placeholder, right=-text_placeholder.left,
                                     bottom=text_placeholder.top + Inches(1))
                # Pic to the right
                position_placeholder(pic, left=text_placeholder.left,
                                     right=-text_placeholder.left,
                                     top=text_placeholder.top + text_placeholder.height + Inches(0.25))
                pic.height = int(pic.width / ar)
                # Check to see if it is off the screen
                distance_over = (pic.height + pic.top - prs.slide_height +
                                 Inches(0.25)) / pic.height
                if distance_over > 0:
                    pic.height = int(pic.height * (1 - distance_over))
                    pic.width = int(pic.width * (1 - distance_over))

                os.remove(output_save_file.format(i + 1))
                text_placeholder.text_frame.paragraphs[0].text = 'Mode {:}'.format(i + 1)
                for label, value in zip(['Frequency: ' + frequency_format + ' Hz', 'Damping: ' + damping_format + ' %'],
                                        [shape.frequency, shape.damping * 100]):
                    p = text_placeholder.text_frame.add_paragraph()
                    p.text = label.format(value)
                    p.level = 1
                text_placeholder.text_frame.fit_text()
                text_placeholder.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    if isinstance(presentation, pptx.presentation.Presentation):
        return prs
    else:
        prs.save(presentation)
        return prs
